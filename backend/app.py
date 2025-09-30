import json
import base64
import zipfile
import io
import difflib
import pdfkit
import logging
import tempfile
import requests 
import os
import threading
import subprocess 
import shutil
import uuid
import zipfile
from datetime import datetime
from typing import List
import subprocess
import tempfile
from docx2pdf import convert
import pypandoc
import PyPDF2
import fitz  # PyMuPDF
import pandas as pd
import pytesseract
from PIL import Image
from PyPDF2 import PdfMerger
from docx import Document as DocxDocument
from docx2pdf import convert as docx2pdf
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pdf2docx import Converter
from pdf2image import convert_from_path
from pptx import Presentation
from werkzeug.utils import secure_filename
from fpdf import FPDF
from openai import OpenAI
import tiktoken
from docx import Document
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from docx import Document
from reportlab.lib.pagesizes import LETTER
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import LETTER
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
# Import all PDF processing functions
from pdf_tools.compare_pdf import compare_pdfs
from pdf_tools.compress import compress_pdf
from pdf_tools.crop_image import crop_image
from pdf_tools.crop_pdf import crop_pdf
from pdf_tools.edit_pdf import edit_pdf
from pdf_tools.excel_to_pdf import excel_to_pdf
from pdf_tools.extract_pages import extract_pages_from_pdf, parse_page_range
# from pdf_tools.extract_text import extract_text_from_pdf

from pdf_tools.html_to_pdf import html_to_pdf
from pdf_tools.jpg_to_pdf import jpg_to_pdf
from pdf_tools.ocr_pdf import ocr_pdf
from pdf_tools.optimize import optimize_pdf
# from pdf_tools.organize_pdf import organize_pdf
from pdf_tools.page_numbers import add_page_numbers
from pdf_tools.pdf_security import encrypt_pdf, decrypt_pdf
from pdf_tools.pdf_to_jpg import pdf_to_images
from pdf_tools.pdf_to_pptx import convert_pdf_to_pptx
from pdf_tools.pdf_to_word_excel import pdf_to_word, pdf_to_excel
from pdf_tools.pptx_to_pdf import convert_pptx_to_pdf
from pdf_tools.redact_pdf import redact_pdf
from pdf_tools.remove_pages import remove_pages_from_pdf
from pdf_tools.repair_pdf import repair_pdf
from pdf_tools.rotate import rotate_pdf
from pdf_tools.scan_pdf import scan_pdf
from pdf_tools.sign_pdf import sign_pdf_basic
from pdf_tools.sign_pdf_cert import sign_pdf_with_cert
from pdf_tools.split import split_pdf_pages
from pdf_tools.unlock_pdf import unlock_pdf
from pdf_tools.watermark import add_watermark
# from pdf_tools.word_to_pdf import word_to_pdf
from flask import Flask, request, jsonify, send_file
from pdf2image import convert_from_bytes
from dotenv import load_dotenv
import time



# Configure Flask app
# load_dotenv()
app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": ["http://localhost:5173", "http://127.0.0.1:5173"]}})
app.config['UPLOAD_FOLDER'] = "Uploads"
app.config['UPLOAD_FOLDER'] = tempfile.gettempdir()
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB limit
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s %(levelname)s: %(message)s')
logger = app.logger



POPPLER_PATH = r"C:\Users\takd2\Release-25.07.0-0 (1)\poppler-25.07.0\Library\bin"

# Allowed file extensions
ALLOWED_EXTENSIONS = {'pdf', 'jpg', 'jpeg', 'png', 'docx', 'xlsx', 'xls', 'pptx', 'html', 'txt', 'csv'}
UPLOAD_FOLDER = os.path.join(os.getcwd(), "Uploads")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)  
# Ensure upload folder exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
if not os.access(app.config['UPLOAD_FOLDER'], os.W_OK):
    logger.error(f"Upload folder {app.config['UPLOAD_FOLDER']} is not writable")
    raise PermissionError(f"Upload folder {app.config['UPLOAD_FOLDER']} is not writable")

# Helper Functions
def allowed_file(filename: str) -> bool:
    """Check if file extension is allowed."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def save_file(file, temp_dir: str, prefix: str = "") -> str:
    """Save uploaded file with a unique filename in a temporary directory."""
    filename = secure_filename(file.filename)
    unique_id = uuid.uuid4().hex
    new_filename = f"{unique_id}_{prefix}{filename}" if prefix else f"{unique_id}_{filename}"
    file_path = os.path.join(temp_dir, new_filename)
    file.save(file_path)
    logger.debug(f"Saved file: {file_path}")
    return file_path

def extract_text_from_pdf(pdf_file_path: str) -> str:
    """Extract text from a PDF file using PyMuPDF."""
    try:
        text = ""
        with fitz.open(pdf_file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text
    except Exception as e:
        logger.error(f"Error extracting text from {pdf_file_path}: {str(e)}")
        raise

def text_from_excel(file_path):
    try:
        df = pd.read_excel(file_path)
        text = "\n".join(df.astype(str).agg(" ".join, axis=1))
        return text
    except Exception as e:
        logger.error(f"Error extracting text from Excel {file_path}: {str(e)}")
        raise

def extract_text_from_word(file_path):
    try:
        doc = DocxDocument(file_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return text
    except Exception as e:
        logger.error(f"Error extracting text from Word {file_path}: {str(e)}")
        raise

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
        raise

def extract_text(file_path, file_type):
    try:
        if file_type == "pdf":
            return extract_text_from_pdf(file_path)
        elif file_type in ["xlsx", "xls"]:
            return extract_text_from_excel(file_path)
        elif file_type == "docx":
            return extract_text_from_word(file_path)
        elif file_type == "pptx":
            return extract_text_from_pptx(file_path)
        elif file_type in ["jpg", "png", "jpeg"]:
            return pytesseract.image_to_string(Image.open(file_path))
        elif file_type in ["txt", "html"]:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        elif file_type == "csv":
            df = pd.read_csv(file_path)
            return df.to_string()
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    except Exception as e:
        logger.error(f"Text extraction failed for {file_path}: {str(e)}")
        raise

def convert_pdf_to_docx(pdf_path, output_path):
    try:
        cv = Converter(pdf_path)
        cv.convert(output_path)
        cv.close()
        return output_path
    except Exception as e:
        logger.error(f"PDF to DOCX conversion failed: {str(e)}")
        raise

def convert_pdf_to_image(pdf_path, output_dir):
    try:
        images = convert_from_path(pdf_path)
        output_paths = []
        for i, image in enumerate(images):
            output_path = os.path.join(output_dir, f"page_{i + 1}.jpg")
            image.save(output_path, "JPEG")
            output_paths.append(output_path)
        return output_paths
    except Exception as e:
        logger.error(f"PDF to image conversion failed: {str(e)}")
        raise

def convert_to_pdf(input_path, file_type, output_path):
    try:
        if file_type == "docx":
            docx2pdf(input_path, output_path)
        elif file_type == "pptx":
            return pptx_to_pdf(input_path)
        elif file_type in ["xlsx", "xls"]:
            return excel_to_pdf(input_path)
        elif file_type in ["jpg", "png", "jpeg"]:
            return image_to_pdf(input_path)
        elif file_type == "html":
            with open(input_path, 'r') as f:
                html_content = f.read()
            return html_to_pdf(html_content)
        elif file_type == "txt":
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            with open(input_path, 'r', encoding='utf-8') as f:
                for line in f:
                    pdf.multi_cell(200, 10, line)
            pdf.output(output_path)
        elif file_type == "csv":
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            df = pd.read_csv(input_path)
            for _, row in df.iterrows():
                line = " | ".join(str(val) for val in row)
                pdf.multi_cell(200, 10, line)
            pdf.output(output_path)
        else:
            raise ValueError(f"Unsupported file type for PDF conversion: {file_type}")
        return output_path
    except Exception as e:
        logger.error(f"Conversion to PDF failed: {str(e)}")
        raise

def perform_ocr(file_path, file_type):
    try:
        if file_type == "pdf":
            return ocr_pdf(file_path)
        elif file_type in ["jpg", "png", "jpeg"]:
            return pytesseract.image_to_string(Image.open(file_path))
        else:
            raise ValueError("OCR only supported for PDF, JPG, PNG")
    except Exception as e:
        logger.error(f"OCR failed for {file_path}: {str(e)}")
        raise


# Core LangChain utilities
from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.documents import Document

# Community integrations (external tools, databases, vectorstores, etc.)
import os
from flask import Flask, request, jsonify
from flask_cors import CORS
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
import google.generativeai as genai
# If you are using Google Gemini
import google.generativeai as genai


# 🔑 Load Gemini API Key from .env
from dotenv import load_dotenv
load_dotenv()
genai.configure(api_key=os.getenv("NEW_KEY_FROM_STUDENT_PLAN"))

#  Load embeddings (must match what you used when saving FAISS)
embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

# 🔹 Load FAISS vectorstore from disk
VECTORSTORE_PATH = "faiss_index"
if os.path.exists(VECTORSTORE_PATH):
    vectorstore = FAISS.load_local(VECTORSTORE_PATH, embeddings, allow_dangerous_deserialization=True)
else:
    vectorstore = None   # Will remain None until a PDF is uploaded and processed

chat_model = genai.GenerativeModel("gemini-2.5-pro")

# Extract text from PDF
def extract_pdf_text(pdf_file):
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return text

# Upload PDF & create embeddings
@app.route("/upload-pdf", methods=["POST"])
def upload_pdf():
    global vectorstore

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "Empty filename"}), 400

    pdf_path = os.path.join("uploads", file.filename)
    file.save(pdf_path)

    # Ensure file saved properly
    if not os.path.exists(pdf_path) or os.path.getsize(pdf_path) == 0:
        return jsonify({"error": "File not saved correctly"}), 400

    # Validate PDF first
    try:
        import fitz
        with fitz.open(pdf_path) as doc:
            if doc.page_count == 0:
                return jsonify({"error": "PDF has no pages"}), 400
    except Exception:
        return jsonify({"error": "Invalid or corrupted PDF"}), 400

    # Process with LangChain
    try:
        loader = PyPDFLoader(pdf_path)
        documents = loader.load()
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        texts = text_splitter.split_documents(documents)

        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        vectorstore = FAISS.from_documents(texts, embeddings)
        vectorstore.save_local("faiss_index")

        return jsonify({"message": "PDF uploaded and processed successfully"})
    except Exception as e:
        return jsonify({"error": f"Processing failed: {str(e)}"}), 500



@app.route("/chat", methods=["POST"])
def chat():
    global vectorstore
    data = request.json
    query = data.get("message")

    if not query:
        return jsonify({"error": "No message provided"}), 400

    docs_text = ""

    # ✅ If PDF index exists, use it
    if vectorstore is None and os.path.exists("faiss_index"):
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        vectorstore = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)

    if vectorstore:
        docs = vectorstore.similarity_search(query, k=3)
        docs_text = "\n\n".join([d.page_content for d in docs])

    # ✅ Build prompt
    if docs_text:
        prompt = f"Answer the question based only on the following PDF content:\n\n{docs_text}\n\nQuestion: {query}"
    else:
        prompt = f"General Chat:\n\nQuestion: {query}"

    # ✅ Gemini call
    model = genai.GenerativeModel("gemini-2.5-pro")
    response = model.generate_content(prompt)

    output_text = ""
    if hasattr(response, "text"):
        output_text = response.text
    elif hasattr(response, "candidates") and response.candidates:
        output_text = response.candidates[0].content.parts[0].text
    
    return jsonify({"response": output_text})

# Convert File
@app.route("/convert", methods=["POST", "OPTIONS"])
def convert_file():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    if 'file' not in request.files or 'target_format' not in request.form:
        return jsonify({"success": False, "error": "File and target format required."}), 400

    file = request.files['file']
    target_format = request.form.get("target_format").lower()
    if not allowed_file(file.filename) or target_format not in ['docx', 'xlsx', 'pptx', 'jpg', 'pdf']:
        return jsonify({"success": False, "error": "Invalid file or target format."}), 400

    try:
        with tempfile.TemporaryDirectory(dir=app.config['UPLOAD_FOLDER']) as temp_dir:
            file_path = save_file(file, temp_dir)
            file_type = file.filename.rsplit('.', 1)[1].lower()
            output_filename = f"converted_{secure_filename(file.filename)}.{target_format}"
            output_path = os.path.join(temp_dir, output_filename)

            if file_type == "pdf":
                if target_format == "docx":
                    output_path = pdf_to_word(file_path)
                elif target_format == "xlsx":
                    output_path = pdf_to_excel(file_path)
                elif target_format == "pptx":
                    output_path = pdf_to_pptx(file_path)
                elif target_format == "jpg":
                    output_paths = pdf_to_jpg(file_path)
                    output_path = output_paths[0]
            elif file_type in ["docx", "xlsx", "xls", "pptx", "jpg", "png", "jpeg", "html", "txt", "csv"]:
                if target_format == "pdf":
                    output_path = convert_to_pdf(file_path, file_type, output_path)
                else:
                    return jsonify({"success": False, "error": "Unsupported conversion."}), 400
            else:
                return jsonify({"success": False, "error": "Unsupported file type."}), 400

            logger.debug(f"Converted {file.filename} to {target_format}")
            return send_file(output_path, download_name=output_filename, as_attachment=True)
    except Exception as e:
        logger.error(f"Conversion failed for {file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Conversion failed: {str(e)}"}), 500


# Merge PDFs
@app.route("/merge", methods=["POST"])
def merge_pdfs():
    try:
        files = request.files.getlist("files")
        if not files:
            return jsonify({"error": "No files uploaded"}), 400

        merger = PdfMerger()
        temp_dir = os.path.join(UPLOAD_FOLDER, f"tmp_{uuid.uuid4().hex}")
        os.makedirs(temp_dir, exist_ok=True)  # ✅ ensure temp directory exists

        output_path = os.path.join(temp_dir, f"merged_{uuid.uuid4()}.pdf")

        for f in files:
            file_path = os.path.join(temp_dir, f.filename)
            f.save(file_path)
            merger.append(file_path)

        merger.write(output_path)
        merger.close()

        return send_file(output_path, as_attachment=True)

    except Exception as e:
        print(f"[ERROR] Merging failed: {e}")
        return jsonify({"error": str(e)}), 500
# Split PDF


@app.route("/split", methods=["POST", "OPTIONS"])
def split_pdf():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    file = request.files.get("file")
    pages = request.form.get("pages")
    if not file or not pages:
        return jsonify({"success": False, "error": "PDF file and pages parameter are required."}), 400
    if not allowed_file(file.filename):
        return jsonify({"success": False, "error": "File type not allowed", "allowed": list(ALLOWED_EXTENSIONS)}), 400

    # create unique temp dir inside Uploads/
    temp_dir = os.path.join(app.config['UPLOAD_FOLDER'], f"temp_{uuid.uuid4().hex}")
    os.makedirs(temp_dir, exist_ok=True)

    try:
        # Save input file into this temp folder
        input_path = save_file(file, temp_dir, "split_")

        # Perform split
        output_path = split_pdf_pages(input_path, pages, temp_dir)

        logger.debug(f"Split PDF {file.filename} into {output_path}")
        return send_file(output_path, as_attachment=True, download_name=f"split_{file.filename}")

    except Exception as e:
        logger.error(f"Splitting failed for {file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Splitting failed: {str(e)}"}), 500

    finally:
        # cleanup temp directory after use
        shutil.rmtree(temp_dir, ignore_errors=True)

# Page Count
@app.route("/page-count", methods=["POST", "OPTIONS"])
def page_count():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    file = request.files.get("file")
    if not file or not allowed_file(file.filename):
        return jsonify({"success": False, "error": "Valid PDF file required."}), 400

    try:
        with tempfile.TemporaryDirectory(dir=app.config['UPLOAD_FOLDER']) as temp_dir:
            input_path = save_file(file, temp_dir, "page_count_")
            with open(input_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                page_count = len(reader.pages)

            return jsonify({"success": True, "pageCount": page_count})
    except Exception as e:
        logger.error(f"Page count failed: {str(e)}")
        return jsonify({"success": False, "error": f"Failed to get page count: {str(e)}"}), 400

# Remove Pages
@app.route("/remove-pages", methods=["POST"])
def remove_pages():
    if "file" not in request.files or "pages" not in request.form:
        return jsonify({"error": "Missing file or pages parameter"}), 400

    file = request.files["file"]
    pages_to_remove = request.form["pages"].strip()

    if not file or not pages_to_remove:
        return jsonify({"error": "Invalid file or pages"}), 400

    try:
        reader = PdfReader(file)
        total_pages = len(reader.pages)
        print(f"📄 PDF has {total_pages} pages")
        print(f"🗑 Pages requested to remove: {pages_to_remove}")

        # Parse pages string -> list of integers
        pages_to_remove_set = set()
        for part in pages_to_remove.split(","):
            part = part.strip()
            if not part:
                continue
            if "-" in part:
                start, end = part.split("-")
                pages_to_remove_set.update(range(int(start), int(end) + 1))
            else:
                pages_to_remove_set.add(int(part))

        # Validate pages (ignore invalid numbers)
        valid_pages = {p for p in pages_to_remove_set if 1 <= p <= total_pages}
        print(f"✅ Valid pages to remove: {valid_pages}")

        writer = PdfWriter()

        # Add pages except those in valid_pages
        for i in range(total_pages):
            if (i + 1) not in valid_pages:
                writer.add_page(reader.pages[i])

        output = io.BytesIO()
        writer.write(output)
        output.seek(0)

        return send_file(
            output,
            as_attachment=True,
            download_name="removed_pages.pdf",
            mimetype="application/pdf",
        )

    except Exception as e:
        print("❌ Error removing pages:", e)  # <-- Full error in backend console
        return jsonify({"error": str(e)}), 500

# Extract Pages
@app.route("/extract-pages", methods=["POST"])
def extract_pages_route():
    try:
        file = request.files["file"]
        pages_str = request.form.get("pages")

        # Save uploaded file
        temp_dir = tempfile.mkdtemp(dir=UPLOAD_FOLDER)
        file_path = os.path.join(temp_dir, secure_filename(file.filename))
        file.save(file_path)

        # Extract pages
        extracted_path = extract_pages_from_pdf(file_path, pages_str, temp_dir)

        print(f"[INFO] 📦 Sending extracted PDF: {extracted_path}")
        return send_file(extracted_path, as_attachment=True)

    except Exception as e:
        print(f"[ERROR] Page extraction failed: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/rotate", methods=["POST"])
def rotate_pdf_route():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files["file"]
        pages_str = request.form.get("pages")
        angle = request.form.get("angle")

        if not pages_str or not angle:
            return jsonify({"error": "Pages and angle are required"}), 400

        try:
            angle = int(angle)
        except ValueError:
            return jsonify({"error": "Angle must be an integer"}), 400

        # Save uploaded file
        upload_dir = os.path.join("Uploads", f"tmp_{uuid.uuid4().hex[:8]}")
        os.makedirs(upload_dir, exist_ok=True)
        input_path = os.path.join(upload_dir, file.filename)
        file.save(input_path)

        # Rotate PDF
        output_path = rotate_pdf(input_path, pages_str, angle, upload_dir)

        print(f"[INFO] 📦 Sending rotated PDF: {output_path}")
        return send_file(output_path, as_attachment=True)

    except Exception as e:
        print(f"[ERROR] Rotation failed: {e}")
        return jsonify({"error": str(e)}), 500


def create_watermark_from_image(image_path, output_path, width=200, height=200):
    """
    Convert an image into a one-page PDF watermark.
    """
    c = canvas.Canvas(output_path, pagesize=letter)
    img = ImageReader(image_path)
    page_width, page_height = letter

    # Place at center by default
    x = (page_width - width) / 2
    y = (page_height - height) / 2

    c.drawImage(img, x, y, width=width, height=height, mask="auto")
    c.save()



def add_watermark(input_pdf, watermark_path):
    """
    Add watermark (PDF or image) to each page of a PDF.
    """
    watermark_pdf = watermark_path
    if watermark_path.lower().endswith((".png", ".jpg", ".jpeg")):
        # Convert image watermark into a temp PDF
        watermark_pdf = watermark_path + ".pdf"
        create_watermark_from_image(watermark_path, watermark_pdf)

    with open(input_pdf, "rb") as f_input, open(watermark_pdf, "rb") as f_watermark:
        pdf_reader = PyPDF2.PdfReader(f_input)
        watermark_reader = PyPDF2.PdfReader(f_watermark)
        watermark_page = watermark_reader.pages[0]

        pdf_writer = PyPDF2.PdfWriter()

        for page in pdf_reader.pages:
            page.merge_page(watermark_page)  # overlay watermark
            pdf_writer.add_page(page)

        # ✅ always use safe temporary output path
        fd, output_path = tempfile.mkstemp(suffix="_watermarked.pdf", dir="Uploads")
        os.close(fd)
        with open(output_path, "wb") as f_output:
            pdf_writer.write(f_output)

    return output_path


@app.route("/watermark", methods=["POST"])
def watermark_pdf():
    """
    API: Upload PDF + watermark (PDF or image).
    """
    try:
        if "file" not in request.files or "watermark" not in request.files:
            return jsonify({"error": "Missing file or watermark"}), 400

        file = request.files["file"]
        watermark = request.files["watermark"]

        # ✅ Create unique temp dir
        temp_dir = tempfile.mkdtemp(dir="Uploads")

        pdf_path = os.path.join(temp_dir, file.filename)
        watermark_path = os.path.join(temp_dir, watermark.filename)

        file.save(pdf_path)
        watermark.save(watermark_path)

        # Add watermark
        output_path = add_watermark(pdf_path, watermark_path)

        print(f"[INFO] ✅ Sending watermarked PDF: {output_path}")
        return send_file(output_path, as_attachment=True)

    except Exception as e:
        print(f"[ERROR] Watermarking failed: {str(e)}")
        return jsonify({"error": str(e)}), 500

# Add Page Numbers

from werkzeug.utils import secure_filename

@app.route("/page-numbers", methods=["POST"])
def add_page_numbers_route():
    try:
        file = request.files["file"]
        position = request.form.get("position", "bottom-right")
        font_size = int(request.form.get("font_size", 12))
        start_page = int(request.form.get("start_page", 1))
        number_format = request.form.get("number_format", "1")

        # Create a unique temp dir for this upload
        temp_dir = tempfile.mkdtemp(dir="Uploads")

        # Sanitize filename (replace spaces & special chars)
        original_filename = file.filename
        safe_filename = "".join(c if c.isalnum() or c in "._-" else "_" for c in original_filename)

        pdf_path = os.path.join(temp_dir, safe_filename)
        file.save(pdf_path)

        # ✅ output file path (not directory!)
        output_path = os.path.join(temp_dir, f"numbered_{safe_filename}")

        # Call PDF tool
        from pdf_tools.page_numbers import add_page_numbers
        result_path = add_page_numbers(
            pdf_path,
            output_path,         # fixed
            position=position,
            font_size=font_size,
            start_page=start_page,
            number_format=number_format
        )

        return send_file(result_path, as_attachment=True)

    except Exception as e:
        app.logger.error(f"[ERROR] Failed to add page numbers to {file.filename}: {e}")
        return jsonify({"error": str(e)}), 500



# Compress PD
@app.route("/compress", methods=["POST"])
def compress_pdf_route():
    try:
        if "file" not in request.files:
            return jsonify({"success": False, "error": "No file uploaded"}), 400

        pdf_file = request.files["file"]

        if not pdf_file.filename.endswith(".pdf"):
            return jsonify({"success": False, "error": "Only PDF files allowed"}), 400

        # save input temporarily
        input_path = os.path.join(app.config["UPLOAD_FOLDER"], pdf_file.filename)
        pdf_file.save(input_path)

        # output path
        output_path = os.path.join(app.config["UPLOAD_FOLDER"], f"compressed_{pdf_file.filename}")

        # call compression
        from pdf_tools.compress import compress_pdf
        compress_pdf(input_path, output_path, image_quality=50)

        return send_file(output_path, as_attachment=True, download_name=f"compressed_{pdf_file.filename}")

    except Exception as e:
        import traceback
        print("❌ Compression failed:", str(e))
        traceback.print_exc()   # shows full error in terminal
        return jsonify({"success": False, "error": str(e)}), 500



# Optimize PDF
@app.route("/optimize", methods=["POST", "OPTIONS"])
def optimize_pdf_route():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    if "file" not in request.files:
        return jsonify({"success": False, "error": "No PDF file uploaded"}), 400

    pdf_file = request.files["file"]
    if not allowed_file(pdf_file.filename):
        return jsonify({"success": False, "error": "File type not allowed", "allowed": list(ALLOWED_EXTENSIONS)}), 400

    try:
        with tempfile.TemporaryDirectory(dir=app.config['UPLOAD_FOLDER']) as temp_dir:
            pdf_path = save_file(pdf_file, temp_dir, "optimize_")
            output_path = optimize_pdf(pdf_path, temp_dir)

            logger.debug(f"Optimized {pdf_file.filename}")
            return send_file(output_path, as_attachment=True, download_name=f"optimized_{pdf_file.filename}")
    except Exception as e:
        logger.error(f"Optimization failed for {pdf_file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Optimization failed: {str(e)}"}), 500

def delete_file_later(path, delay=10):
    def delete_task():
        time.sleep(delay)
        try:
            os.remove(path)
            print(f"[INFO] Deleted temp file: {path}")
        except Exception as e:
            print(f"[WARN] Could not delete temp file: {e}")
    threading.Thread(target=delete_task, daemon=True).start()


@app.route("/pdf-to-images", methods=["POST"])
def pdf_to_images():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    pdf_file = request.files["file"]
    pdf_bytes = pdf_file.read()

    try:
        poppler_path = r"C:\Users\takd2\Release-25.07.0-0 (1)\poppler-25.07.0\Library\bin"
        print(f"[DEBUG] Using Poppler path: {poppler_path}")

        images = convert_from_bytes(pdf_bytes, poppler_path=poppler_path)
        print(f"[DEBUG] Converted {len(images)} pages")

        # Save to temp
        temp_dir = tempfile.gettempdir()
        zip_path = os.path.join(temp_dir, f"pdf_images_{uuid.uuid4().hex}.zip")
        print(f"[DEBUG] Writing zip to: {zip_path}")

        with zipfile.ZipFile(zip_path, "w") as zipf:
            for i, img in enumerate(images):
                img_path = os.path.join(temp_dir, f"page_{i+1}.jpg")
                img.save(img_path, "JPEG")
                zipf.write(img_path, f"page_{i+1}.jpg")
                os.remove(img_path)

        delete_file_later(zip_path, delay=30)
        return send_file(zip_path, as_attachment=True, download_name="pdf_images.zip")

    except Exception as e:
        import traceback
        print("[ERROR]", traceback.format_exc())   # ✅ Full error in terminal
        return jsonify({"error": str(e)}), 500


@app.route("/pdf-to-word", methods=["POST"])
def pdf_to_word():
    print("📥 Incoming request at /pdf-to-word")

    if "file" not in request.files:
        print("❌ No file part in request")
        return jsonify({"error": "No file uploaded"}), 400

    pdf_file = request.files["file"]

    # Save uploaded PDF
    os.makedirs("uploads", exist_ok=True)
    pdf_path = os.path.join("uploads", pdf_file.filename)
    pdf_file.save(pdf_path)

    # Convert PDF -> Word
    word_path = pdf_path.replace(".pdf", ".docx")
    try:
        cv = Converter(pdf_path)
        cv.convert(word_path, start=0, end=None)  # full document
        cv.close()
    except Exception as e:
        print("❌ Conversion failed:", e)
        return jsonify({"error": str(e)}), 500

    # Send back converted Word file
    return send_file(word_path, as_attachment=True)


@app.route("/pdf-to-excel", methods=["POST", "OPTIONS"])
def convert_pdf_to_excel():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    if "file" not in request.files:
        return jsonify({"success": False, "error": "No PDF file uploaded"}), 400

    file = request.files["file"]

    # Ensure PDF files are allowed
    if not allowed_file(file.filename) or file.filename.rsplit('.', 1)[-1].lower() != "pdf":
        return jsonify({"success": False, "error": "File type not allowed", "allowed": ["pdf"]}), 400

    try:
        with tempfile.TemporaryDirectory(dir=app.config['UPLOAD_FOLDER']) as temp_dir:
            file_path = save_file(file, temp_dir, "pdf_to_excel_")
            excel_path = pdf_to_excel(file_path)

            logger.debug(f"Converted {file.filename} to Excel")
            return send_file(
                excel_path,
                as_attachment=True,
                download_name=f"converted_{file.filename.rsplit('.', 1)[0]}.xlsx"
            )
    except Exception as e:
        logger.error(f"PDF to Excel failed for {file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Conversion failed: {str(e)}"}), 500


# Encrypt PDF
@app.route("/encrypt-pdf", methods=["POST", "OPTIONS"])
def encrypt_pdf_route():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    if "file" not in request.files or "password" not in request.form:
        return jsonify({"success": False, "error": "File and password required"}), 400

    file = request.files["file"]
    password = request.form["password"]
    if not allowed_file(file.filename):
        return jsonify({"success": False, "error": "File type not allowed", "allowed": list(ALLOWED_EXTENSIONS)}), 400

    try:
        with tempfile.TemporaryDirectory(dir=app.config['UPLOAD_FOLDER']) as temp_dir:
            file_path = save_file(file, temp_dir, "encrypt_")
            encrypted_path = encrypt_pdf(file_path, password)

            logger.debug(f"Encrypted {file.filename}")
            return send_file(encrypted_path, as_attachment=True, download_name=f"encrypted_{file.filename}")
    except Exception as e:
        logger.error(f"Encryption failed for {file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Encryption failed: {str(e)}"}), 500

# Decrypt PDF
@app.route("/decrypt-pdf", methods=["POST", "OPTIONS"])
def decrypt_pdf_route():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    if "file" not in request.files or "password" not in request.form:
        return jsonify({"success": False, "error": "File and password required"}), 400

    file = request.files["file"]
    password = request.form["password"]
    if not allowed_file(file.filename):
        return jsonify({"success": False, "error": "File type not allowed", "allowed": list(ALLOWED_EXTENSIONS)}), 400

    try:
        with tempfile.TemporaryDirectory(dir=app.config['UPLOAD_FOLDER']) as temp_dir:
            file_path = save_file(file, temp_dir, "decrypt_")
            decrypted_path = decrypt_pdf(file_path, password)

            logger.debug(f"Decrypted {file.filename}")
            return send_file(decrypted_path, as_attachment=True, download_name=f"decrypted_{file.filename}")
    except Exception as e:
        logger.error(f"Decryption failed for {file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Decryption failed: {str(e)}"}), 500




@app.route("/pdf-to-pptx", methods=["POST", "OPTIONS"])
def pdf_to_pptx():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]

    if file.filename == "":
        return jsonify({"error": "Empty filename"}), 400

    file_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(file_path)

    try:
        pptx_path = convert_pdf_to_pptx(file_path)
        return send_file(pptx_path, as_attachment=True)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# PPTX to PDF
@app.route("/pptx-to-pdf", methods=["POST", "OPTIONS"])
def pptx_to_pdf_api():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    if "file" not in request.files:
        return jsonify({"success": False, "error": "No file uploaded"}), 400

    file = request.files["file"]

    if not allowed_file(file.filename):
        return jsonify({"success": False, "error": "File type not allowed. Only PPTX is supported."}), 400

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx", dir=UPLOAD_FOLDER) as tmp_pptx:
            file.save(tmp_pptx.name)
            pptx_path = tmp_pptx.name

        pdf_path = convert_pptx_to_pdf(pptx_path)

        return send_file(pdf_path, as_attachment=True, download_name=f"{os.path.splitext(file.filename)[0]}.pdf")

    except Exception as e:
        logger.error(f"Conversion failed for {file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Conversion failed: {str(e)}"}), 500



# Word to PDF


# Define folders (must exist!)
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "converted"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# Path to LibreOffice (adjust if needed)
# Windows default:
LIBREOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"
# Linux/macOS default:
# LIBREOFFICE_PATH = "/usr/bin/libreoffice"


def word_to_pdf(input_path, output_folder):
    """
    Convert Word file (doc/docx) to PDF using LibreOffice headless mode.
    """
    try:
        result = subprocess.run([
            LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf",
            "--outdir", output_folder, input_path
        ], capture_output=True, text=True, check=True)

        print("✅ LibreOffice STDOUT:", result.stdout)
        print("⚠️ LibreOffice STDERR:", result.stderr)

        filename = os.path.splitext(os.path.basename(input_path))[0] + ".pdf"
        return os.path.join(output_folder, filename)

    except subprocess.CalledProcessError as e:
        print("❌ LibreOffice failed")
        print("STDOUT:", e.stdout)
        print("STDERR:", e.stderr)
        raise


@app.route("/word-to-pdf", methods=["POST", "OPTIONS"])
def convert_word_to_pdf():
    if "file" not in request.files:
        return jsonify({"error": "No file part in request"}), 400

    file = request.files["file"]

    if file.filename == "":
        return jsonify({"error": "No selected file"}), 400

    if not file.filename.lower().endswith((".doc", ".docx")):
        return jsonify({"error": "Only .doc or .docx files are allowed"}), 400

    # Save uploaded file
    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(input_path)

    try:
        # Convert to PDF
        output_path = word_to_pdf(input_path, OUTPUT_FOLDER)

        # Send PDF back to client
        return send_file(output_path, as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 500



# Excel to PDF
@app.route("/excel-to-pdf", methods=["POST"])
def convert_excel_to_pdf():
    if "file" not in request.files:
        return jsonify({"success": False, "error": "No Excel file uploaded"}), 400

    file = request.files["file"]
    if not allowed_file(file.filename):
        return jsonify({"success": False, "error": "File type not allowed"}), 400

    try:
        # Ensure permanent upload folder exists
        os.makedirs("uploads", exist_ok=True)

        # Save file to uploads with unique name
        filename = f"{uuid.uuid4()}_{file.filename}"
        filepath = os.path.abspath(os.path.join("uploads", filename))
        file.save(filepath)

        # Convert to PDF
        pdf_path = excel_to_pdf(filepath, output_folder="uploads")

        return send_file(pdf_path, as_attachment=True,
                         download_name=f"{file.filename.rsplit('.', 1)[0]}.pdf")

    except Exception as e:
        logger.error(f"Excel to PDF failed for {file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Excel to PDF failed: {str(e)}"}), 500


# Edit PDF
@app.route("/edit-pdf", methods=["POST"])
def edit_pdf():
    try:
        file = request.files.get("file")
        if not file:
            return jsonify({"error": "No PDF file provided"}), 400

        options = request.form.get("options")
        options = eval(options)  # Convert string back to dict

        add_text = options.get("addText", "")
        rotate = options.get("rotate", 0)
        watermark = options.get("watermark", "")
        page_numbers = options.get("pageNumbers", False)

        add_image_file = request.files.get("image")

        pdf_bytes = file.read()
        pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")

        # Add image if provided
        img_reader = None
        if add_image_file:
            image = Image.open(add_image_file)
            img_reader = ImageReader(image)

        for page_number, page in enumerate(pdf_doc):
            # Rotate page
            if rotate in [90, 180, 270]:
                page.set_rotation(rotate)

            # Add text
            if add_text:
                rect = fitz.Rect(50, 50, 550, 100)
                page.insert_textbox(rect, add_text, fontsize=20, fontname="helv", color=(0,0,0))

            # Add watermark
            if watermark:
                rect = fitz.Rect(200, 400, 400, 450)
                page.insert_textbox(rect, watermark, fontsize=40, rotate=45, fontname="helv", color=(0.7,0.7,0.7))

            # Add page numbers
            if page_numbers:
                page.insert_text((500, 800), str(page_number+1), fontsize=12, color=(0,0,0))

            # Add image
            if img_reader:
                rect = fitz.Rect(100, 200, 300, 400)
                page.insert_image(rect, stream=add_image_file.read())

        # Save to BytesIO
        output_pdf = BytesIO()
        pdf_doc.save(output_pdf)
        pdf_doc.close()
        output_pdf.seek(0)

        return send_file(output_pdf, as_attachment=True, download_name="edited.pdf", mimetype="application/pdf")

    except Exception as e:
        print("Error editing PDF:", e)
        return jsonify({"error": "Failed to edit PDF"}), 500

@app.route("/sign-pdf", methods=["POST"])
def sign_pdf():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400

    pdf_file = request.files["file"]
    input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4()}_input.pdf")
    output_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4()}_basic_signed.pdf")
    pdf_file.save(input_path)

    try:
        reader = PdfReader(input_path)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        # Add metadata to mark "basic signed"
        writer.add_metadata({"/SignedBy": "Basic Signature"})

        with open(output_path, "wb") as f:
            writer.write(f)

        return send_file(output_path, as_attachment=True, download_name="signed.pdf")

    except Exception as e:
        return jsonify({"error": str(e)}), 500


# --- Certificate Signature Route ---
@app.route("/sign-with-cert", methods=["POST"])
def sign_with_cert():
    if "file" not in request.files or "cert" not in request.files:
        return jsonify({"error": "Missing file or certificate"}), 400

    pdf_file = request.files["file"]
    cert_file = request.files["cert"]
    password = request.form.get("password")

    if not password:
        return jsonify({"error": "Password required for certificate"}), 400

    input_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4()}_input.pdf")
    cert_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4()}_cert.p12")
    output_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4()}_cert_signed.pdf")

    pdf_file.save(input_path)
    cert_file.save(cert_path)

    try:
        signer = signers.SimpleSigner.load_pkcs12(cert_path, password.encode())
        meta = PdfSignatureMetadata(field_name="Signature1")

        with open(input_path, "rb") as inf, open(output_path, "wb") as outf:
            PdfSigner(meta, signer=signer).sign_pdf(inf, output=outf)

        return send_file(output_path, as_attachment=True, download_name="signed.pdf")

    except Exception as e:
        return jsonify({"error": str(e)}), 500
PDFKIT_CONFIG = pdfkit.configuration(wkhtmltopdf=r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe")

@app.route("/html-to-pdf", methods=["POST", "OPTIONS"])
def html_to_pdf():
    # Handle CORS preflight
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    try:
        # Determine input type
        html_content = ""
        if "url" in request.form and request.form["url"]:
            html_content = pdfkit.from_url(request.form["url"], False, configuration=PDFKIT_CONFIG)
        elif "html" in request.form and request.form["html"]:
            html_content = pdfkit.from_string(request.form["html"], False, configuration=PDFKIT_CONFIG)
        else:
            return jsonify({"success": False, "error": "No HTML content or URL provided"}), 400

        # Save PDF to a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
            tmp_file.write(html_content)
            tmp_file_path = tmp_file.name

        # Send PDF to client
        return send_file(tmp_file_path, as_attachment=True, download_name="converted.pdf")

    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

    finally:
        # Cleanup temporary file after response is sent
        if 'tmp_file_path' in locals() and os.path.exists(tmp_file_path):
            try:
                os.remove(tmp_file_path)
            except Exception:
                pass

# Unlock PDF
@app.route("/unlock-pdf", methods=["POST", "OPTIONS"])
def unlock_pdf_route():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    if 'file' not in request.files or "password" not in request.form:
        return jsonify({"success": False, "error": "PDF file and password required"}), 400

    file = request.files['file']
    password = request.form.get("password")
    if not allowed_file(file.filename):
        return jsonify({"success": False, "error": "File type not allowed", "allowed": list(ALLOWED_EXTENSIONS)}), 400

    try:
        with tempfile.TemporaryDirectory(dir=app.config['UPLOAD_FOLDER']) as temp_dir:
            input_path = save_file(file, temp_dir, "unlock_")
            output_path = unlock_pdf(input_path, password)

            logger.debug(f"Unlocked {file.filename}")
            return send_file(output_path, as_attachment=True, download_name=f"unlocked_{file.filename}")
    except Exception as e:
        logger.error(f"Unlocking failed for {file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Unlocking failed: {str(e)}"}), 500


def parse_page_order(page_order_str, total_pages):
    """
    Converts a string like "1,3-5,2" into a list of zero-based page indices.
    """
    if not page_order_str.strip():
        return list(range(total_pages))  # keep original order if empty

    pages = []
    for part in page_order_str.split(','):
        part = part.strip()
        if '-' in part:
            start, end = map(int, part.split('-'))
            pages.extend(range(start-1, end))  # convert to zero-based index
        else:
            pages.append(int(part)-1)
    # Remove duplicates and filter invalid pages
    pages = [p for p in pages if 0 <= p < total_pages]
    return pages

@app.route('/organize-pdf', methods=['POST'])
def organize_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    page_order_str = request.form.get('pageOrder', '')

    try:
        # Read the uploaded PDF
        reader = PdfReader(file)
        writer = PdfWriter()
        total_pages = len(reader.pages)

        # Parse the page order
        pages_to_keep = parse_page_order(page_order_str, total_pages)

        # Reorder pages
        for page_num in pages_to_keep:
            writer.add_page(reader.pages[page_num])

        # Save to a BytesIO object
        output_stream = io.BytesIO()
        writer.write(output_stream)
        output_stream.seek(0)

        return send_file(
            output_stream,
            as_attachment=True,
            download_name="organized.pdf",
            mimetype="application/pdf"
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500
# Repair PDF
@app.route("/repair-pdf", methods=["POST", "OPTIONS"])
def repair_pdf_route():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    if 'file' not in request.files:
        return jsonify({"success": False, "error": "No file uploaded"}), 400

    file = request.files['file']
    if not allowed_file(file.filename):
        return jsonify({"success": False, "error": "File type not allowed", "allowed": list(ALLOWED_EXTENSIONS)}), 400

    try:
        with tempfile.TemporaryDirectory(dir=app.config['UPLOAD_FOLDER']) as temp_dir:
            input_path = save_file(file, temp_dir, "repair_")
            output_path = repair_pdf(input_path)

            logger.debug(f"Repaired {file.filename}")
            return send_file(output_path, as_attachment=True, download_name=f"repaired_{file.filename}")
    except Exception as e:
        logger.error(f"Repair failed for {file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Repair failed: {str(e)}"}), 500

def extract_text(pdf_path):
    text = ""
    with open(pdf_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text
# Compare PDFs
def compute_similarity(text1, text2):
    s = difflib.SequenceMatcher(None, text1, text2)
    return round(s.ratio() * 100, 2)

@app.route("/compare-pdf", methods=["POST"])
def compare_pdf():
    if "file1" not in request.files or "file2" not in request.files:
        return jsonify({"success": False, "error": "Both PDF files are required"}), 400

    file1 = request.files["file1"]
    file2 = request.files["file2"]

    # Save uploaded files properly
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as f1:
        file1.save(f1.name)
        file1_path = f1.name

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as f2:
        file2.save(f2.name)
        file2_path = f2.name

    # Now read after files are fully written
    text1 = extract_text(file1_path)
    text2 = extract_text(file2_path)

    # Compute differences line by line
    diff = list(difflib.unified_diff(text1.splitlines(), text2.splitlines(), lineterm=""))

    # Compute similarity percentage
    similarity = compute_similarity(text1, text2)

    response = {
        "success": True,
        "file1Pages": len(PyPDF2.PdfReader(file1_path).pages),
        "file2Pages": len(PyPDF2.PdfReader(file2_path).pages),
        "differences": diff[:50],  # Limit output
        "similarity": similarity
    }

    return jsonify(response)

# Redact PDF
@app.route("/redact-pdf", methods=["POST"])
def redact_pdf():
    try:
        file = request.files.get("file")
        keywords = request.form.get("keywords", "")

        if not file or not keywords.strip():
            return jsonify({"error": "PDF file and keywords are required"}), 400

        # Save uploaded file
        temp_dir = tempfile.mkdtemp(dir=app.config["UPLOAD_FOLDER"])
        input_path = os.path.join(temp_dir, file.filename)
        file.save(input_path)

        output_path = os.path.join(
            temp_dir, f"redacted_{os.path.splitext(file.filename)[0]}.pdf"
        )

        # Process with PyMuPDF
        doc = fitz.open(input_path)
        words_to_redact = [k.strip() for k in keywords.split(",") if k.strip()]

        for page in doc:
            for word in words_to_redact:
                areas = page.search_for(word)
                for rect in areas:
                    # Add redaction annotation
                    page.add_redact_annot(rect, fill=(0, 0, 0))  # black box

            page.apply_redactions()

        doc.save(output_path)
        doc.close()

        return send_file(output_path, as_attachment=True, download_name="redacted.pdf")

    except Exception as e:
        print(f"[ERROR] Redaction failed: {e}")
        return jsonify({"error": str(e)}), 500




# OCR PDF
@app.route("/ocr-pdf", methods=["POST"])
def ocr_pdf_route():
    try:
        file = request.files["file"]
        input_dir = "Uploads"
        os.makedirs(input_dir, exist_ok=True)
        pdf_path = os.path.join(input_dir, file.filename)
        file.save(pdf_path)

        os.makedirs("Outputs", exist_ok=True)
        base_name = os.path.splitext(file.filename)[0]

        # Run OCR → text
        text_output = os.path.join("Outputs", f"{base_name}_ocr.txt")
        pdf_output = os.path.join("Outputs", f"{base_name}_searchable.pdf")

        result_text = ocr_pdf(pdf_path, text_output)  # your updated function
        # Optional: make searchable PDF (overlay OCR text on original)
        # generate_searchable_pdf(pdf_path, text_output, pdf_output)

        # Decide what to send
        return send_file(text_output, mimetype="text/plain")

    except Exception as e:
        app.logger.error(f"OCR failed for {file.filename}: {e}")
        return jsonify({"error": str(e)}), 500




# Crop PDF
@app.route("/crop-pdf", methods=["POST", "OPTIONS"])
def crop_pdf_route():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No PDF file provided"}), 400

        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "Empty filename"}), 400

        # --- Accept cropBox (array string) OR x0,y0,x1,y1 (separate) ---
        crop_box = None
        if request.form.get("cropBox"):
            try:
                crop_box = tuple(map(float, request.form.get("cropBox").strip("[]").split(",")))
                if len(crop_box) != 4:
                    raise ValueError
            except Exception:
                return jsonify({"error": "Invalid cropBox format"}), 400
        else:
            x0 = request.form.get("x0")
            y0 = request.form.get("y0")
            x1 = request.form.get("x1")
            y1 = request.form.get("y1")
            if not all([x0, y0, x1, y1]):
                return jsonify({"error": "Missing crop coordinates"}), 400
            try:
                crop_box = (float(x0), float(y0), float(x1), float(y1))
            except Exception:
                return jsonify({"error": "Invalid crop coordinates"}), 400

        # Save uploaded file
        filename = secure_filename(file.filename)
        temp_dir = tempfile.mkdtemp(dir=UPLOAD_FOLDER)
        input_path = os.path.join(temp_dir, filename)
        file.save(input_path)

        # Process crop
        output_path = crop_pdf(input_path, crop_box, temp_dir)

        return send_file(output_path, as_attachment=True, download_name=f"cropped_{filename}")

    except PDFProcessingError as e:
        return jsonify({"error": str(e)}), 500
    except Exception as e:
        return jsonify({"error": f"Unexpected error: {str(e)}"}), 500


# Crop Image
@app.route("/crop-image", methods=["POST", "OPTIONS"])
def crop_image_route():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    if "file" not in request.files:
        return jsonify({"success": False, "error": "No image file uploaded"}), 400

    file = request.files["file"]
    if not allowed_file(file.filename):
        return jsonify({"success": False, "error": "File type not allowed", "allowed": list(ALLOWED_EXTENSIONS)}), 400

    try:
        crop_data = request.form
        try:
            left = int(crop_data.get("left", 0))
            upper = int(crop_data.get("top", 0))
            right = int(crop_data.get("right", 300))
            lower = int(crop_data.get("bottom", 300))
            if left < 0 or upper < 0 or right <= left or lower <= upper:
                return jsonify({"success": False, "error": "Invalid crop coordinates"}), 400
        except (ValueError, TypeError):
            return jsonify({"success": False, "error": "Crop coordinates must be numeric"}), 400

        with tempfile.TemporaryDirectory(dir=app.config['UPLOAD_FOLDER']) as temp_dir:
            filepath = save_file(file, temp_dir, "crop_image_")
            output_path = crop_image(filepath, (left, upper, right, lower))

            logger.debug(f"Cropped image {file.filename}")
            return send_file(output_path, as_attachment=True, download_name=f"cropped_{file.filename}")
    except Exception as e:
        logger.error(f"Image cropping failed for {file.filename}: {str(e)}")
        return jsonify({"success": False, "error": f"Cropping failed: {str(e)}"}), 500


@app.route("/jpg-to-pdf", methods=["POST"])
def jpg_to_pdf_route():
    if "files" not in request.files:
        return jsonify({"error": "No files uploaded"}), 400

    # 👇 get multiple uploaded files
    files = request.files.getlist("files")
    if not files:
        return jsonify({"error": "Empty file list"}), 400

    image_paths = []
    for file in files:
        filename = secure_filename(file.filename)
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(file_path)
        image_paths.append(file_path)  # ✅ add each image path

    # Convert to one PDF
    pdf_path = jpg_to_pdf(image_paths)

    # Return merged PDF
    return send_file(pdf_path, as_attachment=True, download_name="converted.pdf")



# Health Check
@app.route('/health', methods=["GET", "OPTIONS"])
def health_check():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "GET, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    checks = {
        'flask_running': True,
        'upload_dir_exists': os.path.exists(app.config['UPLOAD_FOLDER']),
        'upload_dir_writable': os.access(app.config['UPLOAD_FOLDER'], os.W_OK),
        'disk_space': shutil.disk_usage(app.config['UPLOAD_FOLDER']).free > 100 * 1024 * 1024
    }
    status = "healthy" if all(checks.values()) else "degraded"

    logger.debug(f"Health check: {status}")
    return jsonify({
        "success": True,
        "status": status,
        "version": "1.0.0",
        "timestamp": datetime.utcnow().isoformat(),
        "checks": checks
    }), 200 if status == "healthy" else 500

# Home Endpoint
@app.route("/", methods=["GET", "OPTIONS"])
def home():
    if request.method == "OPTIONS":
        response = jsonify({"success": True})
        response.headers.add("Access-Control-Allow-Methods", "GET, OPTIONS")
        response.headers.add("Access-Control-Allow-Headers", "Content-Type")
        return response

    endpoints = {
        "upload_chat_pdf": {
            "method": "POST",
            "path": "/upload-chat-pdf",
            "description": "Upload a file and extract text for chat",
            "parameters": {"file": "PDF, Excel, Word, PPTX, JPG, PNG, TXT, CSV, HTML file"}
        },
        "api_chat": {
            "method": "POST",
            "path": "/api/chat",
            "description": "General chat with AI",
            "parameters": {"messages": "List of chat messages"}
        },
        "ask_question": {
            "method": "POST",
            "path": "/ask-question",
            "description": "Ask a question about extracted text",
            "parameters": {"text": "Extracted text", "question": "User question"}
        },
        "convert": {
            "method": "POST",
            "path": "/convert",
            "description": "Convert file to PDF, DOCX, XLSX, PPTX, or JPG",
            "parameters": {"file": "Source file", "target_format": "docx, xlsx, pptx, jpg, pdf"}
        },
        "ocr": {
            "method": "POST",
            "path": "/ocr",
            "description": "Perform OCR on PDF or image",
            "parameters": {"file": "PDF, JPG, PNG file"}
        },
        "merge": {
            "method": "POST",
            "path": "/merge",
            "description": "Merge multiple PDF files",
            "parameters": {"files": "PDF files (2+ required)"}
        },
        "split": {
            "method": "POST",
            "path": "/split",
            "description": "Split PDF by pages",
            "parameters": {"file": "PDF file", "pages": "Page ranges (e.g., '1,3,5-7')"}
        },
        "page_count": {
            "method": "POST",
            "path": "/page-count",
            "description": "Get page count of a PDF",
            "parameters": {"file": "PDF file"}
        },
        "remove_pages": {
            "method": "POST",
            "path": "/remove-pages",
            "description": "Remove specified pages from PDF",
            "parameters": {"file": "PDF file", "pages": "Page ranges (e.g., '2,4-5')"}
        },
        "extract_pages": {
            "method": "POST",
            "path": "/extract-pages",
            "description": "Extract specified pages from PDF",
            "parameters": {"file": "PDF file", "pages": "Page ranges (e.g., '1,3-5')"}
        },
        "rotate": {
            "method": "POST",
            "path": "/rotate",
            "description": "Rotate PDF pages",
            "parameters": {"file": "PDF file", "pages": "Page ranges (e.g., '1,3-5')",
                           "angle": "Rotation angle (90, 180, 270)"}
        },
        "watermark": {
            "method": "POST",
            "path": "/watermark",
            "description": "Add watermark to PDF",
            "parameters": {"file": "PDF file", "watermark": "Watermark PDF/image"}
        },
        "add_page_numbers": {
            "method": "POST",
            "path": "/add-page-numbers",
            "description": "Add page numbers to PDF",
            "parameters": {"file": "PDF file"}
        },
        "compress": {
            "method": "POST",
            "path": "/compress",
            "description": "Compress a PDF",
            "parameters": {"file": "PDF file"}
        },
        "optimize": {
            "method": "POST",
            "path": "/optimize",
            "description": "Optimize a PDF",
            "parameters": {"file": "PDF file"}
        },
        "pdf_to_jpg": {
            "method": "POST",
            "path": "/pdf-to-jpg",
            "description": "Convert PDF to JPG images",
            "parameters": {"file": "PDF file"}
        },
        "jpg_to_pdf": {
            "method": "POST",
            "path": "/jpg-to-pdf",
            "description": "Convert JPG images to PDF",
            "parameters": {"files": "JPG files"}
        },
        "pdf_to_word": {
            "method": "POST",
            "path": "/pdf-to-word",
            "description": "Convert PDF to Word",
            "parameters": {"file": "PDF file"}
        },
        "pdf_to_excel": {
            "method": "POST",
            "path": "/pdf-to-excel",
            "description": "Convert PDF to Excel",
            "parameters": {"file": "PDF file"}
        },
        "encrypt_pdf": {
            "method": "POST",
            "path": "/encrypt-pdf",
            "description": "Encrypt a PDF with a password",
            "parameters": {"file": "PDF file", "password": "Encryption password"}
        },
        "decrypt_pdf": {
            "method": "POST",
            "path": "/decrypt-pdf",
            "description": "Decrypt a PDF with a password",
            "parameters": {"file": "PDF file", "password": "Decryption password"}
        },
        "pdf_to_pptx": {
            "method": "POST",
            "path": "/pdf-to-pptx",
            "description": "Convert PDF to PowerPoint",
            "parameters": {"file": "PDF file"}
        },
        "pptx_to_pdf": {
            "method": "POST",
            "path": "/pptx-to-pdf",
            "description": "Convert PowerPoint to PDF",
            "parameters": {"file": "PPTX file"}
        },
        "word_to_pdf": {
            "method": "POST",
            "path": "/word-to-pdf",
            "description": "Convert Word to PDF",
            "parameters": {"file": "DOCX file"}
        },
        "excel_to_pdf": {
            "method": "POST",
            "path": "/excel-to-pdf",
            "description": "Convert Excel to PDF",
            "parameters": {"file": "XLSX or XLS file"}
        },
        "edit_pdf": {
            "method": "POST",
            "path": "/edit-pdf",
            "description": "Edit a PDF with specified changes",
            "parameters": {"file": "PDF file", "edits": "JSON object with edit instructions"}
        },
        "sign_pdf": {
            "method": "POST",
            "path": "/sign-pdf",
            "description": "Sign a PDF with a basic signature",
            "parameters": {"file": "PDF file"}
        },
        "sign_with_cert": {
            "method": "POST",
            "path": "/sign-with-cert",
            "description": "Sign a PDF with a digital certificate",
            "parameters": {"file": "PDF file", "cert": "Certificate file", "password": "Certificate password"}
        },
        "html_to_pdf": {
            "method": "POST",
            "path": "/html-to-pdf",
            "description": "Convert HTML content to PDF",
            "parameters": {"html": "HTML content string"}
        },
        "unlock_pdf": {
            "method": "POST",
            "path": "/unlock-pdf",
            "description": "Unlock a password-protected PDF",
            "parameters": {"file": "PDF file", "password": "Decryption password"}
        },
        "organize_pdf": {
            "method": "POST",
            "path": "/organize-pdf",
            "description": "Rearrange PDF pages in specified order",
            "parameters": {"file": "PDF file", "page_order": "Comma-separated list of page numbers"}
        },
        "repair_pdf": {
            "method": "POST",
            "path": "/repair-pdf",
            "description": "Repair a corrupted PDF file",
            "parameters": {"file": "PDF file"}
        },
        "scan_pdf": {
            "method": "POST",
            "path": "/scan-pdf",
            "description": "Convert an image to a searchable PDF",
            "parameters": {"file": "Image file (JPG, PNG)"}
        },
        "compare_pdf": {
            "method": "POST",
            "path": "/compare-pdf",
            "description": "Compare two PDF files and highlight differences",
            "parameters": {"files": "Exactly 2 PDF files"}
        },
        "redact_pdf": {
            "method": "POST",
            "path": "/redact-pdf",
            "description": "Redact specified keywords from a PDF",
            "parameters": {"file": "PDF file", "keywords": "Comma-separated list of keywords to redact"}
        },
        "ocr_pdf": {
            "method": "POST",
            "path": "/ocr-pdf",
            "description": "Perform OCR on a PDF to extract text",
            "parameters": {"file": "PDF file"}
        },
        "crop_pdf": {
            "method": "POST",
            "path": "/crop-pdf",
            "description": "Crop a PDF to specified dimensions",
            "parameters": {"file": "PDF file", "x0": "Left coordinate", "y0": "Top coordinate",
                           "x1": "Right coordinate",
                           "y1": "Bottom coordinate"}
        },
        "crop_image": {
            "method": "POST",
            "path": "/crop-image",
            "description": "Crop an image to specified dimensions",
            "parameters": {"file": "Image file (JPG, PNG)", "left": "Left coordinate", "top": "Top coordinate",
                           "right": "Right coordinate", "bottom": "Bottom coordinate"}
        },
        "health": {
            "method": "GET",
            "path": "/health",
            "description": "Check the health status of the API",
            "parameters": {}
        }

    }

    logger.debug("Home endpoint accessed")
    return jsonify({
        "success": True,
        "message": "Welcome to the PDF Processing API",
        "endpoints": endpoints,
        "version": "1.0.0",
        "timestamp": datetime.utcnow().isoformat()
    }), 200
    

if __name__ == '__main__':
    logger.info("Starting AiPdf Backend")
    app.run(debug=True, host='0.0.0.0', port=5000)