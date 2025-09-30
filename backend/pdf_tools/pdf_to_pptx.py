import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image

PPTX_OUTPUT_FOLDER = "pptx_output"
os.makedirs(PPTX_OUTPUT_FOLDER, exist_ok=True)

def convert_pdf_to_pptx(pdf_path):
    try:
        # Open the PDF
        pdf_document = fitz.open(pdf_path)

        # Create a blank presentation
        prs = Presentation()
        prs.slide_width = Inches(11.69)   # A4 landscape
        prs.slide_height = Inches(8.27)

        blank_slide_layout = prs.slide_layouts[6]  # empty slide

        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]

            # Render page to an image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # higher res
            img_bytes = pix.tobytes("png")

            # Save image temporarily
            img = Image.open(BytesIO(img_bytes))
            img_path = os.path.join(PPTX_OUTPUT_FOLDER, f"page_{page_num+1}.png")
            img.save(img_path)

            # Add slide with image
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(
                img_path,
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        # Save PPTX file
        base = os.path.basename(pdf_path)
        name, _ = os.path.splitext(base)
        pptx_filename = f"{name}.pptx"
        pptx_path = os.path.join(PPTX_OUTPUT_FOLDER, pptx_filename)
        prs.save(pptx_path)

        return pptx_path

    except Exception as e:
        print("Error in convert_pdf_to_pptx:", str(e))
        raise
